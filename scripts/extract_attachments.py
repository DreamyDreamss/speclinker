# STATUS: 완료
#!/usr/bin/env python3
"""
extract_attachments.py — 변경 첨부(PPT/Word/Excel/PDF) 텍스트 추출 (선택 의존성 lazy)
Usage: python extract_attachments.py <attach_dir> [out_md]
출력: out_md (기본 attach_dir/../_extracted.md)
"""
import os, sys, glob

def _pptx(p):
    try:
        from pptx import Presentation
    except ImportError:
        return None
    out = []
    for i, s in enumerate(Presentation(p).slides, 1):
        txt = [sh.text for sh in s.shapes if getattr(sh, 'has_text_frame', False) and sh.text.strip()]
        if txt:
            out.append(f'[슬라이드 {i}]\n' + '\n'.join(txt))
    return '\n\n'.join(out)

def _docx(p):
    try:
        import docx
    except ImportError:
        return None
    return '\n'.join(par.text for par in docx.Document(p).paragraphs if par.text.strip())

def _xlsx(p):
    try:
        import openpyxl
    except ImportError:
        return None
    wb = openpyxl.load_workbook(p, read_only=True, data_only=True)
    out = []
    for ws in wb.worksheets:
        out.append(f'[시트 {ws.title}]')
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                out.append('\t'.join(cells))
    return '\n'.join(out)

def _txt(p):
    return open(p, encoding='utf-8', errors='ignore').read()

EXT = {'.pptx': _pptx, '.docx': _docx, '.xlsx': _xlsx,
       '.txt': _txt, '.md': _txt, '.csv': _txt}

def main():
    if len(sys.argv) < 2:
        print('Usage: extract_attachments.py <attach_dir> [out_md]')
        return 1
    adir = sys.argv[1]
    out = sys.argv[2] if len(sys.argv) > 2 else os.path.join(os.path.dirname(adir.rstrip('/\\')), '_extracted.md')
    blocks = []
    for fp in sorted(glob.glob(os.path.join(adir, '*'))):
        ext = os.path.splitext(fp)[1].lower()
        name = os.path.basename(fp)
        if ext == '.pdf':
            blocks.append(f'## {name}\n[PDF — 에이전트가 Read 도구로 직접 읽기 권장]')
            continue
        fn = EXT.get(ext)
        if not fn:
            blocks.append(f'## {name}\n[추출 미지원 포맷 {ext} — 내용 직접 입력 요청]')
            continue
        txt = fn(fp)
        if txt is None:
            blocks.append(f'## {name}\n[추출 불가 — 라이브러리 미설치(pptx/docx/openpyxl). 직접 입력 요청]')
        else:
            blocks.append(f'## {name}\n{txt.strip()}')
    open(out, 'w', encoding='utf-8').write('\n\n---\n\n'.join(blocks) + '\n')
    print(f'추출 완료: {out} ({len(blocks)}개 파일)')
    return 0

if __name__ == '__main__':
    sys.exit(main())
